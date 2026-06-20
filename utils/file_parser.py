"""
Extracts plain text from uploaded course documents so it can be stored,
searched, and sent to the AI summarizer.
"""
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


class ExtractionError(Exception):
    pass


def extract_text(filepath, file_type):
    file_type = file_type.lower()
    try:
        if file_type == "pdf":
            return _extract_pdf(filepath)
        if file_type == "docx":
            return _extract_docx(filepath)
        if file_type == "pptx":
            return _extract_pptx(filepath)
    except Exception as exc:  # noqa: BLE001 - we want to surface any parser failure
        raise ExtractionError(f"Could not read {file_type.upper()} file: {exc}") from exc

    raise ExtractionError(f"Unsupported file type: {file_type}")


def _extract_pdf(filepath):
    reader = PdfReader(filepath)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages).strip()


def _extract_docx(filepath):
    doc = DocxDocument(filepath)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts).strip()


def _extract_pptx(filepath):
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
    return "\n\n".join(parts).strip()
